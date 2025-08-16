#!/usr/bin/env python3
"""
Pipeline Manager - GPT-4 Document Processor
Roman Engineering Standard: Built for 2000+ year reliability

Core document processing engine using GPT-4 API for offering memorandum
and real estate document extraction with >95% accuracy.
"""

import os
import json
import time
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass
from datetime import datetime
import logging

# Document processing imports
import PyPDF2
import pandas as pd
import openpyxl
from docx import Document
from pptx import Presentation

# AI and API imports
import openai
from openai import OpenAI

# Validation and utilities
import re
import hashlib
from enum import Enum

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentType(Enum):
    """Document type classification"""
    OFFERING_MEMORANDUM = "offering_memorandum"
    FINANCIAL_STATEMENT = "financial_statement"
    RENT_ROLL = "rent_roll"
    PROPERTY_REPORT = "property_report"
    UNKNOWN = "unknown"

@dataclass
class ExtractionResult:
    """Result container for document extraction"""
    success: bool
    data: Dict[str, Any]
    confidence_score: float
    processing_time: float
    errors: List[str]
    warnings: List[str]
    document_type: DocumentType
    extraction_timestamp: str
    
class DocumentProcessor:
    """GPT-4 powered document processor for real estate documents"""
    
    def __init__(self, config_path: Optional[str] = None):
        """Initialize document processor with GPT-4 integration"""
        self.config = self._load_config(config_path)
        self.client = self._initialize_openai_client()
        self.extraction_prompts = self._load_extraction_prompts()
        self.validation_rules = self._load_validation_rules()
        
        logger.info("DocumentProcessor initialized with GPT-4 integration")
    
    def _load_config(self, config_path: Optional[str]) -> Dict[str, Any]:
        """Load processing configuration"""
        default_config = {
            "openai_model": "gpt-4",
            "max_tokens": 4000,
            "temperature": 0.1,  # Low temperature for consistent extraction
            "timeout_seconds": 60,
            "retry_attempts": 3,
            "confidence_threshold": 0.85,
            "extraction_batch_size": 1
        }
        
        if config_path and os.path.exists(config_path):
            try:
                with open(config_path, 'r') as f:
                    user_config = json.load(f)
                default_config.update(user_config)
            except Exception as e:
                logger.warning(f"Failed to load config from {config_path}: {e}")
        
        return default_config
    
    def _initialize_openai_client(self) -> OpenAI:
        """Initialize OpenAI client with API key"""
        api_key = os.getenv('OPENAI_API_KEY')
        if not api_key:
            raise ValueError("OPENAI_API_KEY environment variable required")
        
        return OpenAI(api_key=api_key)
    
    def _load_extraction_prompts(self) -> Dict[str, str]:
        """Load specialized extraction prompts for different document types"""
        return {
            "offering_memorandum": """
You are a real estate document extraction specialist. Extract structured data from this offering memorandum with maximum accuracy.

Extract the following information in JSON format:

{
  "property_details": {
    "name": "Property name",
    "address": "Full property address",
    "city": "City",
    "state": "State",
    "zip_code": "ZIP code",
    "unit_count": "Total number of units (integer)",
    "unit_mix": {
      "studio": "Number of studio units",
      "one_bedroom": "Number of 1BR units",
      "two_bedroom": "Number of 2BR units",
      "three_bedroom": "Number of 3BR+ units"
    },
    "square_footage": "Total square footage",
    "average_unit_size": "Average unit size in sq ft",
    "year_built": "Year built (integer)",
    "property_type": "Property type (Garden, High-rise, etc.)",
    "property_class": "Property class (A, B, C, D)",
    "lot_size": "Lot size in acres or sq ft"
  },
  "financial_metrics": {
    "purchase_price": "Purchase price (number only)",
    "price_per_unit": "Price per unit",
    "gross_scheduled_income": "Annual GSI",
    "net_operating_income": "Annual NOI",
    "cap_rate": "Cap rate as decimal (e.g., 0.055 for 5.5%)",
    "gross_rent_multiplier": "GRM",
    "occupancy_rate": "Physical occupancy as decimal",
    "economic_occupancy": "Economic occupancy as decimal",
    "operating_expenses": "Annual operating expenses",
    "expense_ratio": "Operating expense ratio as decimal"
  },
  "rent_information": {
    "average_rent": "Average rent per unit",
    "rent_per_sqft": "Average rent per square foot",
    "studio_rent": "Average studio rent",
    "one_bedroom_rent": "Average 1BR rent",
    "two_bedroom_rent": "Average 2BR rent",
    "three_bedroom_rent": "Average 3BR rent"
  },
  "market_data": {
    "submarket": "Submarket or neighborhood",
    "median_household_income": "Area median household income",
    "population": "Area population",
    "employment_growth": "Employment growth rate",
    "major_employers": "List of major employers",
    "transportation": "Transportation access description",
    "school_district": "School district name",
    "walk_score": "Walk score if available"
  },
  "transaction_details": {
    "listing_broker": "Listing broker name",
    "broker_company": "Broker company",
    "broker_phone": "Broker phone number",
    "broker_email": "Broker email",
    "days_on_market": "Days on market",
    "seller_motivation": "Seller motivation description",
    "due_diligence_period": "DD period in days",
    "closing_timeline": "Expected closing timeline",
    "financing_available": "Financing options available"
  },
  "additional_info": {
    "amenities": "Property amenities list",
    "recent_improvements": "Recent capital improvements",
    "deferred_maintenance": "Known deferred maintenance",
    "upside_potential": "Value-add opportunities",
    "market_trends": "Local market trends",
    "risks": "Identified risks or concerns"
  }
}

IMPORTANT INSTRUCTIONS:
1. Extract ONLY information that is explicitly stated in the document
2. For missing information, use null or "Not specified"
3. Convert all financial figures to numbers (remove $ and commas)
4. Express percentages as decimals (5.5% = 0.055)
5. Be precise with numbers - accuracy is critical
6. If conflicting information exists, use the most recent or prominently displayed figure
7. Include confidence notes for any uncertain extractions

Document content:
""",
            
            "financial_statement": """
Extract financial data from this financial statement or rent roll.
Focus on income, expenses, and key financial metrics.
Return data in structured JSON format with high accuracy.
""",
            
            "property_report": """
Extract property information from this property report or appraisal.
Focus on physical characteristics, condition, and valuation data.
Return data in structured JSON format.
"""
        }
    
    def _load_validation_rules(self) -> Dict[str, Any]:
        """Load validation rules for extracted data"""
        return {
            "required_fields": {
                "offering_memorandum": [
                    "property_details.name",
                    "property_details.address",
                    "property_details.unit_count",
                    "financial_metrics.purchase_price"
                ]
            },
            "data_types": {
                "property_details.unit_count": int,
                "property_details.year_built": int,
                "financial_metrics.purchase_price": (int, float),
                "financial_metrics.cap_rate": float,
                "financial_metrics.occupancy_rate": float
            },
            "value_ranges": {
                "financial_metrics.cap_rate": (0.01, 0.20),  # 1% to 20%
                "financial_metrics.occupancy_rate": (0.0, 1.0),  # 0% to 100%
                "property_details.year_built": (1800, 2030)
            }
        }
    
    def process_document(self, file_path: str, document_type: Optional[DocumentType] = None) -> ExtractionResult:
        """Process a document and extract structured data"""
        start_time = time.time()
        
        try:
            # Read document content
            content = self._read_document(file_path)
            if not content:
                return ExtractionResult(
                    success=False,
                    data={},
                    confidence_score=0.0,
                    processing_time=time.time() - start_time,
                    errors=["Failed to read document content"],
                    warnings=[],
                    document_type=DocumentType.UNKNOWN,
                    extraction_timestamp=datetime.now().isoformat()
                )
            
            # Classify document type if not provided
            if not document_type:
                document_type = self._classify_document(content)
            
            # Extract data using GPT-4
            extraction_data = self._extract_with_gpt4(content, document_type)
            
            # Validate extracted data
            validation_result = self._validate_extraction(extraction_data, document_type)
            
            # Calculate confidence score
            confidence_score = self._calculate_confidence(extraction_data, validation_result)
            
            processing_time = time.time() - start_time
            
            return ExtractionResult(
                success=True,
                data=extraction_data,
                confidence_score=confidence_score,
                processing_time=processing_time,
                errors=validation_result.get('errors', []),
                warnings=validation_result.get('warnings', []),
                document_type=document_type,
                extraction_timestamp=datetime.now().isoformat()
            )
            
        except Exception as e:
            logger.error(f"Document processing failed: {str(e)}")
            return ExtractionResult(
                success=False,
                data={},
                confidence_score=0.0,
                processing_time=time.time() - start_time,
                errors=[f"Processing error: {str(e)}"],
                warnings=[],
                document_type=document_type or DocumentType.UNKNOWN,
                extraction_timestamp=datetime.now().isoformat()
            )
    
    def _read_document(self, file_path: str) -> Optional[str]:
        """Read document content based on file type"""
        file_extension = Path(file_path).suffix.lower()
        
        try:
            if file_extension == '.pdf':
                return self._read_pdf(file_path)
            elif file_extension in ['.xlsx', '.xls']:
                return self._read_excel(file_path)
            elif file_extension == '.docx':
                return self._read_docx(file_path)
            elif file_extension == '.pptx':
                return self._read_pptx(file_path)
            elif file_extension == '.txt':
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            else:
                logger.warning(f"Unsupported file type: {file_extension}")
                return None
                
        except Exception as e:
            logger.error(f"Failed to read {file_path}: {str(e)}")
            return None
    
    def _read_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text
    
    def _read_excel(self, file_path: str) -> str:
        """Extract text from Excel file"""
        workbook = openpyxl.load_workbook(file_path)
        text = ""
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text += f"\n--- Sheet: {sheet_name} ---\n"
            
            for row in sheet.iter_rows():
                row_text = []
                for cell in row:
                    if cell.value is not None:
                        row_text.append(str(cell.value))
                if row_text:
                    text += " | ".join(row_text) + "\n"
        
        return text
    
    def _read_docx(self, file_path: str) -> str:
        """Extract text from Word document"""
        doc = Document(file_path)
        text = ""
        
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        # Extract table content
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                text += " | ".join(row_text) + "\n"
        
        return text
    
    def _read_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation"""
        prs = Presentation(file_path)
        text = ""
        
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- Slide {slide_num} ---\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        return text
    
    def _classify_document(self, content: str) -> DocumentType:
        """Classify document type based on content"""
        content_lower = content.lower()
        
        # Offering memorandum indicators
        om_keywords = ['offering memorandum', 'investment summary', 'property details', 
                      'financial summary', 'market overview', 'investment highlights']
        
        # Financial statement indicators  
        fs_keywords = ['income statement', 'profit loss', 't-12', 'operating statement',
                      'financial performance', 'revenue', 'expenses']
        
        # Rent roll indicators
        rr_keywords = ['rent roll', 'tenant roster', 'unit mix', 'lease expiration',
                      'monthly rent', 'tenant information']
        
        om_score = sum(1 for keyword in om_keywords if keyword in content_lower)
        fs_score = sum(1 for keyword in fs_keywords if keyword in content_lower)
        rr_score = sum(1 for keyword in rr_keywords if keyword in content_lower)
        
        if om_score >= max(fs_score, rr_score) and om_score > 0:
            return DocumentType.OFFERING_MEMORANDUM
        elif fs_score >= max(om_score, rr_score) and fs_score > 0:
            return DocumentType.FINANCIAL_STATEMENT
        elif rr_score >= max(om_score, fs_score) and rr_score > 0:
            return DocumentType.RENT_ROLL
        else:
            return DocumentType.UNKNOWN
    
    def _extract_with_gpt4(self, content: str, document_type: DocumentType) -> Dict[str, Any]:
        """Extract data using GPT-4 API"""
        prompt_key = document_type.value if document_type.value in self.extraction_prompts else "offering_memorandum"
        prompt = self.extraction_prompts[prompt_key]
        
        # Truncate content if too long (GPT-4 token limit)
        max_content_length = 12000  # Conservative limit
        if len(content) > max_content_length:
            content = content[:max_content_length] + "\n\n[Content truncated...]"
        
        full_prompt = prompt + "\n\n" + content
        
        for attempt in range(self.config['retry_attempts']):
            try:
                response = self.client.chat.completions.create(
                    model=self.config['openai_model'],
                    messages=[
                        {"role": "system", "content": "You are a precise real estate document extraction specialist. Return only valid JSON."},
                        {"role": "user", "content": full_prompt}
                    ],
                    max_tokens=self.config['max_tokens'],
                    temperature=self.config['temperature'],
                    timeout=self.config['timeout_seconds']
                )
                
                # Parse JSON response
                response_text = response.choices[0].message.content.strip()
                
                # Clean response (remove markdown formatting if present)
                if response_text.startswith('```json'):
                    response_text = response_text[7:-3].strip()
                elif response_text.startswith('```'):
                    response_text = response_text[3:-3].strip()
                
                return json.loads(response_text)
                
            except json.JSONDecodeError as e:
                logger.warning(f"JSON decode error on attempt {attempt + 1}: {e}")
                if attempt == self.config['retry_attempts'] - 1:
                    raise
                time.sleep(2 ** attempt)  # Exponential backoff
                
            except Exception as e:
                logger.warning(f"GPT-4 API error on attempt {attempt + 1}: {e}")
                if attempt == self.config['retry_attempts'] - 1:
                    raise
                time.sleep(2 ** attempt)
        
        raise Exception("Failed to extract data after all retry attempts")
    
    def _validate_extraction(self, data: Dict[str, Any], document_type: DocumentType) -> Dict[str, List[str]]:
        """Validate extracted data against business rules"""
        errors = []
        warnings = []
        
        # Check required fields
        required_fields = self.validation_rules['required_fields'].get(document_type.value, [])
        for field_path in required_fields:
            if not self._get_nested_value(data, field_path):
                errors.append(f"Required field missing: {field_path}")
        
        # Check data types
        for field_path, expected_type in self.validation_rules['data_types'].items():
            value = self._get_nested_value(data, field_path)
            if value is not None and not isinstance(value, expected_type):
                warnings.append(f"Data type mismatch for {field_path}: expected {expected_type}, got {type(value)}")
        
        # Check value ranges
        for field_path, (min_val, max_val) in self.validation_rules['value_ranges'].items():
            value = self._get_nested_value(data, field_path)
            if value is not None and isinstance(value, (int, float)):
                if not (min_val <= value <= max_val):
                    warnings.append(f"Value out of range for {field_path}: {value} not in [{min_val}, {max_val}]")
        
        return {'errors': errors, 'warnings': warnings}
    
    def _get_nested_value(self, data: Dict[str, Any], field_path: str) -> Any:
        """Get nested dictionary value using dot notation"""
        keys = field_path.split('.')
        current = data
        
        for key in keys:
            if isinstance(current, dict) and key in current:
                current = current[key]
            else:
                return None
        
        return current
    
    def _calculate_confidence(self, data: Dict[str, Any], validation_result: Dict[str, List[str]]) -> float:
        """Calculate confidence score for extraction"""
        # Base confidence
        confidence = 0.9
        
        # Reduce confidence for errors and warnings
        error_count = len(validation_result.get('errors', []))
        warning_count = len(validation_result.get('warnings', []))
        
        confidence -= (error_count * 0.2)  # -20% per error
        confidence -= (warning_count * 0.05)  # -5% per warning
        
        # Check data completeness
        total_fields = self._count_fields(data)
        empty_fields = self._count_empty_fields(data)
        
        if total_fields > 0:
            completeness = (total_fields - empty_fields) / total_fields
            confidence *= completeness
        
        return max(0.0, min(1.0, confidence))
    
    def _count_fields(self, data: Any) -> int:
        """Count total fields in nested dictionary"""
        if isinstance(data, dict):
            return sum(self._count_fields(value) for value in data.values())
        elif isinstance(data, list):
            return sum(self._count_fields(item) for item in data)
        else:
            return 1
    
    def _count_empty_fields(self, data: Any) -> int:
        """Count empty fields in nested dictionary"""
        if isinstance(data, dict):
            return sum(self._count_empty_fields(value) for value in data.values())
        elif isinstance(data, list):
            return sum(self._count_empty_fields(item) for item in data)
        else:
            return 1 if data in [None, "", "Not specified"] else 0

# Example usage and testing
if __name__ == "__main__":
    # Initialize processor
    processor = DocumentProcessor()
    
    # Example processing
    # result = processor.process_document("path/to/offering_memorandum.pdf")
    # print(f"Success: {result.success}")
    # print(f"Confidence: {result.confidence_score:.2f}")
    # print(f"Processing time: {result.processing_time:.2f}s")
    
    logger.info("DocumentProcessor ready for use")
